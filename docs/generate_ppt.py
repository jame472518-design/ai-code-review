"""Generate SOP PowerPoint for ai-code-review cross-platform installation & usage."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────
BG_DARK   = RGBColor(0x1E, 0x1E, 0x2E)   # dark background
BG_SLIDE  = RGBColor(0x24, 0x27, 0x3A)   # slide body
ACCENT    = RGBColor(0x89, 0xB4, 0xFA)   # blue accent
ACCENT2   = RGBColor(0xA6, 0xE3, 0xA1)   # green
ACCENT3   = RGBColor(0xFA, 0xB3, 0x87)   # peach/orange
WHITE     = RGBColor(0xCD, 0xD6, 0xF4)   # text white
GRAY      = RGBColor(0x6C, 0x70, 0x86)   # dim text
CODE_BG   = RGBColor(0x31, 0x34, 0x4B)   # code block bg
WIN_CLR   = RGBColor(0x74, 0xC7, 0xEC)   # Windows blue
MAC_CLR   = RGBColor(0xF5, 0xC2, 0xE7)   # macOS pink
LNX_CLR   = RGBColor(0xA6, 0xE3, 0xA1)   # Linux green
ACCENT_DIM = RGBColor(0x45, 0x47, 0x5A)  # subtle separator

# ── Typography ──
FONT_TITLE = "Segoe UI"
FONT_BODY  = "Segoe UI"
FONT_CODE  = "Cascadia Code"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

_slide_counter = [0]  # mutable counter for slide numbers


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, fill_color, corner_radius=Inches(0.15)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if shape.adjustments and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.05
    return shape


def _add_rect(slide, left, top, width, height, fill_color):
    """Add a plain rectangle (no rounded corners)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name=None):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name or FONT_BODY
    p.alignment = alignment
    return p


def _add_para(tf, text, size=16, color=WHITE, bold=False, space_before=Pt(2), space_after=Pt(1), indent=0, font_name=None):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name or FONT_BODY
    p.space_before = space_before
    p.space_after = space_after
    if indent:
        p.level = indent
    return p


def _add_bullet(tf, text, size=15, color=WHITE, indent=0):
    return _add_para(tf, text, size=size, color=color, indent=indent)


def _add_accent_line(slide, left, top, width, color=ACCENT, height=Inches(0.04)):
    """Add a thin horizontal accent line."""
    return _add_rect(slide, left, top, width, height, color)


def _add_slide_number(slide, number):
    """Add slide number at bottom-right."""
    tb = _add_text_box(slide, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35))
    _set_text(tb.text_frame, str(number), size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)


def _add_footer_bar(slide):
    """Add a subtle branded footer bar at bottom."""
    _add_rect(slide, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2), ACCENT_DIM)
    tb = _add_text_box(slide, Inches(0.5), Inches(7.28), Inches(5), Inches(0.25))
    _set_text(tb.text_frame, "AI Code Review  |  Cross-Platform SOP", size=8, color=GRAY)


def _finish_slide(slide):
    """Add common elements (slide number + footer) to every slide."""
    _slide_counter[0] += 1
    _add_slide_number(slide, _slide_counter[0])
    _add_footer_bar(slide)


# ── Helper: title slide ──
def make_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, BG_DARK)

    # Decorative top stripe
    _add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

    # Decorative side accent (left bar)
    _add_rect(slide, Inches(0.6), Inches(1.8), Inches(0.06), Inches(4.0), ACCENT_DIM)

    # Title
    tb = _add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.5))
    _set_text(tb.text_frame, title, size=48, color=ACCENT, bold=True,
              alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    # Accent line under title
    _add_accent_line(slide, Inches(4.5), Inches(3.5), Inches(4.3), ACCENT)

    if subtitle:
        tb2 = _add_text_box(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(1))
        _set_text(tb2.text_frame, subtitle, size=22, color=GRAY,
                  alignment=PP_ALIGN.CENTER, font_name=FONT_BODY)

    _finish_slide(slide)
    return slide


# ── Helper: section title ──
def make_section_slide(number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_DARK)

    # Decorative horizontal lines
    _add_accent_line(slide, Inches(2), Inches(2.6), Inches(3.5), ACCENT_DIM)
    _add_accent_line(slide, Inches(7.8), Inches(2.6), Inches(3.5), ACCENT_DIM)

    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.9), Inches(1.85), Inches(1.5), Inches(1.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    tf = circ.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = _set_text(tf, str(number), size=42, color=BG_DARK, bold=True,
                  alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    # Title
    tb = _add_text_box(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(1.2))
    _set_text(tb.text_frame, title, size=36, color=WHITE, bold=True,
              alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    # Subtitle accent line
    _add_accent_line(slide, Inches(5.5), Inches(5.0), Inches(2.3), ACCENT)

    _finish_slide(slide)
    return slide


# ── Helper: 3-column platform slide ──
def make_3col_slide(title, win_lines, mac_lines, linux_lines, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_DARK)

    # Title bar
    tb = _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    _set_text(tb.text_frame, title, size=28, color=ACCENT, bold=True, font_name=FONT_TITLE)

    # Accent line under title
    _add_accent_line(slide, Inches(0.5), Inches(0.95), Inches(12.3), ACCENT_DIM)

    col_w = Inches(3.9)
    col_h = Inches(5.2)
    gap = Inches(0.3)
    left_start = Inches(0.4)
    top = Inches(1.2)

    platforms = [
        ("Windows", WIN_CLR, win_lines),
        ("macOS",   MAC_CLR, mac_lines),
        ("Linux (Ubuntu)", LNX_CLR, linux_lines),
    ]

    for i, (pname, pcolor, lines) in enumerate(platforms):
        left = left_start + i * (col_w + gap)

        # Card background
        _add_shape(slide, left, top, col_w, col_h, CODE_BG)

        # Colored top stripe
        _add_rect(slide, left + Inches(0.1), top + Inches(0.02), col_w - Inches(0.2), Inches(0.05), pcolor)

        # Platform label
        label = _add_text_box(slide, left + Inches(0.15), top + Inches(0.15), col_w - Inches(0.3), Inches(0.5))
        _set_text(label.text_frame, pname, size=18, color=pcolor, bold=True, font_name=FONT_TITLE)

        # Content
        content = _add_text_box(slide, left + Inches(0.15), top + Inches(0.6), col_w - Inches(0.3), col_h - Inches(0.75))
        tf = content.text_frame
        tf.word_wrap = True
        first = True
        for line in lines:
            if first:
                _set_text(tf, line, size=13, color=WHITE, font_name=FONT_BODY)
                first = False
            else:
                is_cmd = line.startswith("$ ") or line.startswith("# ")
                clr = ACCENT2 if is_cmd else WHITE
                sz = 12 if is_cmd else 13
                fn = FONT_CODE if is_cmd else FONT_BODY
                _add_para(tf, line, size=sz, color=clr, space_before=Pt(2), space_after=Pt(1), font_name=fn)

    # Note at bottom
    if note:
        nb = _add_text_box(slide, Inches(0.5), Inches(6.6), Inches(12), Inches(0.6))
        _set_text(nb.text_frame, note, size=13, color=GRAY)

    _finish_slide(slide)
    return slide


# ── Helper: content slide ──
def make_content_slide(title, bullets, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_DARK)

    # Title
    tb = _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    _set_text(tb.text_frame, title, size=28, color=ACCENT, bold=True, font_name=FONT_TITLE)

    # Accent line under title
    _add_accent_line(slide, Inches(0.5), Inches(0.95), Inches(12.3), ACCENT_DIM)

    # Body card
    _add_shape(slide, Inches(0.4), Inches(1.1), Inches(12.5), Inches(5.5), CODE_BG)

    # Left accent bar on card
    _add_rect(slide, Inches(0.4), Inches(1.25), Inches(0.06), Inches(5.2), ACCENT)

    body = _add_text_box(slide, Inches(0.75), Inches(1.3), Inches(11.9), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True

    first = True
    for item in bullets:
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}

        sz = opts.get("size", 16)
        clr = opts.get("color", WHITE)
        b = opts.get("bold", False)
        ind = opts.get("indent", 0)
        is_cmd = text.strip().startswith("$ ") or text.strip().startswith("# ")
        fn = FONT_CODE if is_cmd and sz <= 15 else FONT_BODY

        if first:
            _set_text(tf, text, size=sz, color=clr, bold=b, font_name=fn)
            first = False
        else:
            _add_para(tf, text, size=sz, color=clr, bold=b, indent=ind,
                      space_before=Pt(opts.get("space_before", 3)),
                      space_after=Pt(1), font_name=fn)

    if note:
        nb = _add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5))
        _set_text(nb.text_frame, note, size=12, color=GRAY)

    _finish_slide(slide)
    return slide


# ── Helper: table slide ──
def make_table_slide(title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_DARK)

    # Title
    tb = _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    _set_text(tb.text_frame, title, size=28, color=ACCENT, bold=True, font_name=FONT_TITLE)

    # Accent line under title
    _add_accent_line(slide, Inches(0.5), Inches(0.95), Inches(12.3), ACCENT_DIM)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_w = Inches(12)
    tbl_h = Inches(min(0.5 * n_rows + 0.5, 5.5))
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.65), Inches(1.2), tbl_w, tbl_h)
    table = table_shape.table

    # Header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = BG_DARK
            p.font.name = FONT_TITLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT

    # Rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = WHITE
                p.font.name = FONT_BODY
            cell.fill.solid()
            cell.fill.fore_color.rgb = CODE_BG if i % 2 == 0 else BG_SLIDE

    _finish_slide(slide)
    return slide


# ════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════

# ── Slide 1: Cover ──
make_title_slide(
    "AI Code Review",
    "跨平台安裝與使用 SOP  —  Windows / macOS / Linux"
)

# ── Slide 2: TOC ──
make_content_slide("目錄", [
    ("1.  環境需求", {"size": 20, "bold": True, "color": ACCENT}),
    ("2.  安裝 Ollama（本地 LLM）", {"size": 20, "bold": True, "color": ACCENT}),
    ("3.  安裝 ai-review CLI", {"size": 20, "bold": True, "color": ACCENT}),
    ("4.  設定 LLM Provider", {"size": 20, "bold": True, "color": ACCENT}),
    ("5.  安裝 Git Hooks", {"size": 20, "bold": True, "color": ACCENT}),
    ("6.  日常使用流程", {"size": 20, "bold": True, "color": ACCENT}),
    ("7.  互動式 Commit 助手", {"size": 20, "bold": True, "color": ACCENT}),
    ("8.  使用情境", {"size": 20, "bold": True, "color": ACCENT}),
    ("9.  更換/設定 LLM Model", {"size": 20, "bold": True, "color": ACCENT}),
    ("10. 常用指令速查表", {"size": 20, "bold": True, "color": ACCENT}),
    ("11. 疑難排解", {"size": 20, "bold": True, "color": ACCENT}),
    ("12. 完整移除", {"size": 20, "bold": True, "color": ACCENT}),
])

# ── Slide 3: Prerequisites table ──
make_table_slide("1. 環境需求",
    ["項目", "版本", "說明"],
    [
        ["Python", ">= 3.10", "核心執行環境"],
        ["Git", "任意版本", "版本管理 + hooks"],
        ["pip", "任意版本", "Python 套件管理"],
        ["Ollama（選用）", "任意版本", "本地 LLM 推理引擎"],
    ]
)

# ── Slide 4: Prerequisites per platform ──
make_3col_slide("1. 環境需求 — 安裝方式",
    [  # Windows
        "確認版本：",
        "$ python --version",
        "$ git --version",
        "",
        "安裝方式：",
        "• python.org 下載安裝",
        "  （勾選 Add to PATH）",
        "• 或使用 Anaconda",
        "",
        "注意：",
        "• 建議使用 Git Bash",
        "• CMD 用雙引號 commit",
    ],
    [  # macOS
        "確認版本：",
        "$ python3 --version",
        "$ git --version",
        "",
        "安裝方式：",
        "$ brew install python@3.12",
        "$ brew install git",
        "",
        "macOS 通常自帶 Git",
    ],
    [  # Linux
        "確認版本：",
        "$ python3 --version",
        "$ git --version",
        "",
        "安裝方式：",
        "$ sudo apt update",
        "$ sudo apt install python3 \\",
        "    python3-pip python3-venv git",
    ],
)

# ── Slide 5: Section - Ollama ──
make_section_slide(2, "安裝 Ollama（本地 LLM）")

# ── Slide 6: Install Ollama ──
make_3col_slide("2. 安裝 Ollama",
    [  # Windows
        "1. 下載安裝檔：",
        "   ollama.com/download/windows",
        "",
        "2. 執行安裝程式",
        "   （自動啟動服務）",
        "",
        "3. 開啟 CMD/PowerShell：",
        "$ ollama --version",
        "$ ollama pull llama3.2",
        "$ ollama list",
    ],
    [  # macOS
        "方式一：Homebrew",
        "$ brew install ollama",
        "",
        "方式二：",
        "   ollama.com/download/mac",
        "",
        "啟動服務：",
        "$ ollama serve &",
        "",
        "下載模型：",
        "$ ollama pull llama3.2",
        "$ ollama list",
    ],
    [  # Linux
        "一鍵安裝：",
        "$ curl -fsSL \\",
        "  https://ollama.com/install.sh | sh",
        "",
        "$ ollama --version",
        "$ ollama pull llama3.2",
        "",
        "確認服務：",
        "$ systemctl status ollama",
        "$ sudo systemctl start ollama",
        "$ sudo systemctl enable ollama",
    ],
)

# ── Slide 7: Model comparison ──
make_table_slide("2. 模型推薦",
    ["模型", "大小", "速度", "品質", "適用情境"],
    [
        ["llama3.2", "3.2B", "~8 秒", "中", "日常開發（推薦）"],
        ["llama3.1", "8B", "~15 秒", "中高", "品質優先"],
        ["codellama", "7B", "~4 分鐘", "高", "品質最好但最慢"],
    ]
)

# ── Slide 8: Section - Install CLI ──
make_section_slide(3, "安裝 ai-review CLI")

# ── Slide 9: Install CLI ──
make_3col_slide("3. 安裝 ai-review CLI",
    [  # Windows
        "從原始碼安裝：",
        "$ git clone https://github.com/",
        "    jame472518-design/ai-code-review.git",
        "$ cd ai-code-review",
        "$ pip install -e .",
        "",
        "驗證：",
        "$ ai-review --help",
        "",
        "若找不到指令，確認 PATH：",
        "Anaconda: C:\\Users\\<USER>\\",
        "  anaconda3\\Scripts\\",
    ],
    [  # macOS
        "從原始碼安裝：",
        "$ git clone https://github.com/",
        "    jame472518-design/ai-code-review.git",
        "$ cd ai-code-review",
        "$ python3 -m venv .venv",
        "$ source .venv/bin/activate",
        "$ pip install -e .",
        "",
        "驗證：",
        "$ ai-review --help",
    ],
    [  # Linux
        "從原始碼安裝：",
        "$ git clone https://github.com/",
        "    jame472518-design/ai-code-review.git",
        "$ cd ai-code-review",
        "$ python3 -m venv .venv",
        "$ source .venv/bin/activate",
        "$ pip install -e .",
        "",
        "驗證：",
        "$ ai-review --help",
    ],
    note="或直接安裝：pip install git+https://github.com/jame472518-design/ai-code-review.git"
)

# ── Slide 10: Section - Config ──
make_section_slide(4, "設定 LLM Provider")

# ── Slide 11: Provider config ──
make_content_slide("4. 設定 LLM Provider（全平台通用）", [
    ("方案 A：本地 Ollama（推薦，無需外網）", {"size": 18, "bold": True, "color": ACCENT2}),
    ("$ ai-review config set provider default ollama", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("$ ai-review config set ollama base_url http://localhost:11434", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("$ ai-review config set ollama model llama3.2", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("", {"size": 8}),
    ("方案 B：OpenAI（精確度最高）", {"size": 18, "bold": True, "color": ACCENT3}),
    ("$ ai-review config set provider default openai", {"size": 14, "color": ACCENT3, "indent": 1}),
    ("$ ai-review config set openai api_key_env OPENAI_API_KEY", {"size": 14, "color": ACCENT3, "indent": 1}),
    ("$ ai-review config set openai model gpt-4o", {"size": 14, "color": ACCENT3, "indent": 1}),
    ("", {"size": 8}),
    ("方案 C：企業內部 LLM", {"size": 18, "bold": True, "color": WIN_CLR}),
    ("$ ai-review config set provider default enterprise", {"size": 14, "color": WIN_CLR, "indent": 1}),
    ("$ ai-review config set enterprise base_url https://llm.company.com", {"size": 14, "color": WIN_CLR, "indent": 1}),
    ("", {"size": 8}),
    ("驗證連線：ai-review health-check", {"size": 16, "bold": True, "color": WHITE}),
], note="設定檔位置：~/.config/ai-code-review/config.toml")

# ── Slide 12: Commit config ──
make_content_slide("4. 設定 Commit 相關（全平台通用）", [
    ("設定專案 ID（自動前綴 commit message）", {"size": 18, "bold": True, "color": WHITE}),
    ("$ ai-review config set commit project_id \"PROJ-1\"", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("", {"size": 8}),
    ("初始化 commit 模板", {"size": 18, "bold": True, "color": WHITE}),
    ("$ ai-review config init-template", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("模板檔會複製到 ~/.config/ai-code-review/commit-template.txt", {"size": 14, "color": GRAY, "indent": 1}),
    ("", {"size": 8}),
    ("初始化 AI 生成 prompt", {"size": 18, "bold": True, "color": WHITE}),
    ("$ ai-review config init-generate-prompt", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("prompt 檔會複製到 ~/.config/ai-code-review/generate-prompt.txt", {"size": 14, "color": GRAY, "indent": 1}),
    ("", {"size": 8}),
    ("限縮審查副檔名（預設：不限制，審查所有檔案）", {"size": 18, "bold": True, "color": WHITE}),
    ("$ ai-review config set review include_extensions \"c,cpp,h,py\"", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("預設：空（全部檔案），有需要時再設定限縮", {"size": 14, "color": GRAY, "indent": 1}),
    ("", {"size": 8}),
    ("驗證設定", {"size": 18, "bold": True, "color": WHITE}),
    ("$ ai-review config show", {"size": 14, "color": ACCENT2, "indent": 1}),
])

# ── Slide 13: Section - Hooks ──
make_section_slide(5, "安裝 Git Hooks")

# ── Slide 14: Hook installation ──
make_content_slide("5. 安裝 Git Hooks（全平台通用）", [
    ("方式一：Template Hooks（推薦）", {"size": 18, "bold": True, "color": ACCENT2}),
    ("$ ai-review hook install --template    # 一次性全域設定", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("$ cd /path/to/repo", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("$ ai-review hook enable                # 自動設定 + 安裝 hooks", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("", {"size": 8}),
    ("確認狀態：ai-review hook status", {"size": 16, "bold": True, "color": WHITE}),
    ("", {"size": 8}),
    ("多個 Repo 的批量管理 → 見「情境二」（P22-P23）", {"size": 16, "color": GRAY}),
], note="hook enable 會同時設定 git config + 複製 hook 腳本到 .git/hooks/")

# ── Slide 15: Hook types ──
make_table_slide("5. Hook 類型說明",
    ["Hook", "觸發時機", "功能", "阻擋？"],
    [
        ["pre-commit", "git commit", "AI Code Review 程式碼", "critical/error 時阻擋"],
        ["prepare-commit-msg", "git commit", "互動式 Commit 助手", "否"],
        ["commit-msg", "git commit", "檢查 message 格式", "格式錯誤時阻擋"],
        ["pre-push", "git push", "AI Review 所有待推 commits", "critical/error 時阻擋"],
    ]
)

# ── Slide 16: Section - Daily Use ──
make_section_slide(6, "日常使用流程")

# ── Slide 17: Daily workflow ──
make_content_slide("6. 日常使用流程", [
    ("⚠ 首次使用？先完成以下設定（做過可跳過）", {"size": 16, "bold": True, "color": ACCENT3}),
    ("$ ai-review hook install --template    # 一次性全域設定", {"size": 13, "color": GRAY, "indent": 1}),
    ("$ ai-review hook enable                # 在 repo 裡啟用", {"size": 13, "color": GRAY, "indent": 1}),
    ("", {"size": 6}),
    ("正常 Commit 流程", {"size": 20, "bold": True, "color": ACCENT}),
    ("", {"size": 6}),
    ("1. 修改程式碼", {"size": 16, "color": WHITE}),
    ("2. git add <files>                         # Stage 變更", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("3. git commit                              # 觸發 hooks", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("   → prepare-commit-msg：互動式選單", {"size": 14, "color": GRAY, "indent": 2}),
    ("   → 編輯器開啟，確認 message", {"size": 14, "color": GRAY, "indent": 2}),
    ("   → commit-msg：格式檢查", {"size": 14, "color": GRAY, "indent": 2}),
    ("4. git push origin main                    # 觸發 pre-push", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("", {"size": 8}),
    ("直接 Commit（跳過互動選單）", {"size": 20, "bold": True, "color": ACCENT}),
    ("git commit -m \"[fix] resolve camera crash\"", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("", {"size": 8}),
    ("暫時跳過所有 Hooks", {"size": 20, "bold": True, "color": ACCENT3}),
    ("git commit --no-verify -m \"[hotfix] emergency\"", {"size": 14, "color": ACCENT3, "indent": 1}),
])

# ── Slide 18: Section - Interactive ──
make_section_slide(7, "互動式 Commit 助手")

# ── Slide 19: Interactive menu ──
make_content_slide("7. 互動式 Commit 助手", [
    ("git commit（不帶 -m）時自動出現選單：", {"size": 18, "color": WHITE}),
    ("", {"size": 6}),
    ("  Commit Message Assistant", {"size": 16, "bold": True, "color": ACCENT}),
    ("    1  Load template       - 載入模板", {"size": 15, "color": WIN_CLR}),
    ("    2  Manual draft        - 輸入草稿 → AI 優化", {"size": 15, "color": MAC_CLR}),
    ("    3  Fixed questions     - 固定問題 → 回答 → AI 生成", {"size": 15, "color": ACCENT}),
    ("    4  LLM interview       - AI 問你問題 → 生成", {"size": 15, "color": ACCENT3}),
    ("    5  LLM auto-generate   - AI 從 diff 自動生成", {"size": 15, "color": LNX_CLR}),
    ("    s  Skip                - 跳過，直接進編輯器", {"size": 15, "color": GRAY}),
    ("", {"size": 10}),
    ("選項 1：載入 commit template 到編輯器", {"size": 16, "color": WIN_CLR}),
    ("選項 2：輸入草稿 → AI 讓描述更專業精確 → 編輯器確認", {"size": 16, "color": MAC_CLR}),
    ("選項 3：用固定問題檔提問 → 回答後 AI 生成 → 編輯器確認", {"size": 16, "color": ACCENT}),
    ("選項 4：AI 問你中英雙語問題 → 根據答案生成 → 編輯器確認", {"size": 16, "color": ACCENT3}),
    ("選項 5：兩階段自動生成（分析 diff → 生成 message）→ 編輯器確認", {"size": 16, "color": LNX_CLR}),
    ("選項 s：直接進入編輯器手動編寫", {"size": 16, "color": GRAY}),
    ("", {"size": 8}),
    ("所有選項最終都進入編輯器，讓你確認後才真正 commit", {"size": 16, "bold": True, "color": ACCENT2}),
])

# ── Slide 20: Section - Scenarios ──
make_section_slide(8, "使用情境")

# ── Slide 21: Scenario - Single dev ──
make_content_slide("8. 情境一：個人開發者", [
    ("一次性設定", {"size": 18, "bold": True, "color": ACCENT}),
    ("$ pip install -e .", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("$ ai-review config set provider default ollama", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("$ ai-review config set ollama model llama3.2", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("$ ai-review config init-template", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("$ ai-review config init-generate-prompt", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("$ ai-review hook enable", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("", {"size": 8}),
    ("日常使用 — hooks 自動運作", {"size": 18, "bold": True, "color": ACCENT}),
    ("$ git add main.py", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("$ git commit          → 互動選單自動出現", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("                      → 編輯器確認 message", {"size": 14, "color": GRAY, "indent": 2}),
    ("                      → 格式自動檢查", {"size": 14, "color": GRAY, "indent": 2}),
    ("$ git push            → pre-push AI 審查", {"size": 14, "color": ACCENT2, "indent": 1}),
])

# ── Slide 22: Scenario - Multi repo (daily) ──
make_content_slide("8. 情境二：管理多個 Repo — 日常使用", [
    ("⚠ 首次使用？先完成全域設定（做過可跳過）", {"size": 16, "bold": True, "color": ACCENT3}),
    ("$ ai-review hook install --template    # 一次性安裝 template hooks", {"size": 13, "color": GRAY, "indent": 1}),
    ("", {"size": 6}),
    ("先預覽所有 repo 狀態", {"size": 18, "bold": True, "color": ACCENT}),
    ("$ ai-review hook enable --all /workspace --list", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("  camera-hal: disabled    audio-hal: disabled", {"size": 13, "color": GRAY, "indent": 1}),
    ("  kernel: disabled        framework: disabled", {"size": 13, "color": GRAY, "indent": 1}),
    ("", {"size": 6}),
    ("方法 A：Android BSP 團隊（repo forall）", {"size": 18, "bold": True, "color": ACCENT2}),
    ("$ repo forall -c 'ai-review hook enable'", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("  透過 repo forall 對所有 manifest 內的 repo 逐一啟用", {"size": 13, "color": GRAY, "indent": 1}),
    ("", {"size": 6}),
    ("方法 B：通用（掃描目錄下所有 git repo）", {"size": 18, "bold": True, "color": ACCENT3}),
    ("$ ai-review hook enable --all /workspace", {"size": 14, "color": ACCENT3, "indent": 1}),
    ("  自動掃描 /workspace 下所有 git repo 並啟用", {"size": 13, "color": GRAY, "indent": 1}),
], note="hook enable 會同時設定 git config + 複製 hook 腳本到 .git/hooks/")

# ── Slide: Scenario - Commit choices ──
make_table_slide("8. 情境三：Commit 互動選單怎麼選",
    ["選項", "適合什麼時候", "說明"],
    [
        ["1 Load template", "忘了格式", "載入模板填寫"],
        ["2 Manual draft", "有草稿想優化", "輸入草稿 → AI 讓描述更專業"],
        ["3 Fixed questions", "團隊標準化", "固定問題檔 → 回答 → AI 生成"],
        ["4 LLM interview", "不知怎麼寫", "AI 問你問題 → 根據答案生成"],
        ["5 LLM auto-generate", "懶得寫", "兩階段：分析 diff → 自動生成"],
        ["s Skip", "自己寫", "直接進編輯器"],
    ]
)

# ── Slide 25: Section - Model Config ──
make_section_slide(9, "更換/設定 LLM Model")

# ── Slide 22: Switch model ──
make_content_slide("9. 切換 Model（全平台通用）", [
    ("下載新模型", {"size": 18, "bold": True, "color": ACCENT}),
    ("$ ollama pull llama3.2", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("$ ollama pull llama3.1", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("$ ollama list                              # 查看已下載的模型", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("", {"size": 6}),
    ("切換 ai-review 使用的 Model", {"size": 18, "bold": True, "color": ACCENT}),
    ("$ ai-review config set ollama model llama3.2", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("$ ai-review health-check                   # 驗證", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("", {"size": 6}),
    ("臨時切換（不改設定，僅本次生效）", {"size": 18, "bold": True, "color": ACCENT}),
    ("$ ai-review --model codellama", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("  同 provider，只換模型（例：Ollama 換用 codellama）", {"size": 13, "color": GRAY, "indent": 1}),
    ("$ ai-review --provider openai --model gpt-4o", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("  連 provider 一起換（例：從 Ollama 切到 OpenAI）", {"size": 13, "color": GRAY, "indent": 1}),
], note="小模型（3B）速度快但可能誤報 ｜ 大模型（7B+）品質好但較慢")

# ── Slide 23: Model comparison ──
make_table_slide("9. 模型比較",
    ["模型", "參數量", "速度", "品質", "適用情境"],
    [
        ["llama3.2", "3.2B", "~8 秒", "中", "日常開發（推薦）"],
        ["llama3.1", "8B", "~15 秒", "中高", "品質優先"],
        ["codellama", "7B", "~4 分鐘", "高", "程式碼專用，品質最好"],
        ["qwen2.5-coder", "7B", "~20 秒", "高", "程式碼專用替代方案"],
    ]
)

# ── Slide 24: Remote Ollama architecture ──
make_content_slide("9. 遠端 Ollama Server 架構", [
    ("多人共用一台 LLM Server（同網域內）", {"size": 18, "bold": True, "color": ACCENT}),
    ("", {"size": 6}),
    ("  使用者 A  ──→", {"size": 15, "color": WIN_CLR}),
    ("  使用者 B  ──→  LLM Server（192.168.1.100）", {"size": 15, "color": MAC_CLR}),
    ("  使用者 C  ──→  Ollama + llama3.2 + GPU", {"size": 15, "color": LNX_CLR}),
    ("", {"size": 4}),
    ("  diff 會傳送到 LLM Server 進行推理，使用者不需安裝 GPU", {"size": 13, "color": GRAY}),
    ("  ※ 建議使用內網 Server，避免程式碼外洩", {"size": 13, "color": ACCENT3}),
    ("", {"size": 8}),
    ("每位使用者只需設定一次（指向 Server IP）", {"size": 18, "bold": True, "color": ACCENT2}),
    ("$ ai-review config set ollama base_url http://192.168.1.100:11434", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("$ ai-review health-check", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("", {"size": 6}),
    ("切回本地 Ollama", {"size": 16, "bold": True, "color": ACCENT3}),
    ("$ ai-review config set ollama base_url http://localhost:11434", {"size": 14, "color": ACCENT3, "indent": 1}),
])

# ── Slide 25: Remote server setup per platform ──
make_3col_slide("9. Server 端設定 — 讓 Ollama 監聽所有網卡",
    [  # Windows
        "系統環境變數（永久）：",
        "  設定 > 系統 > 進階系統設定",
        "  > 環境變數",
        "  新增系統變數：",
        "  OLLAMA_HOST = 0.0.0.0",
        "  重啟 Ollama",
        "",
        "PowerShell（暫時）：",
        "$ $env:OLLAMA_HOST=\"0.0.0.0\"",
        "$ ollama serve",
        "",
        "防火牆：",
        "  輸入規則 > TCP 11434",
    ],
    [  # macOS
        "設定環境變數：",
        "$ launchctl setenv \\",
        "    OLLAMA_HOST \"0.0.0.0\"",
        "",
        "重啟 Ollama App",
        "",
        "或終端啟動：",
        "$ OLLAMA_HOST=0.0.0.0 \\",
        "    ollama serve",
        "",
        "防火牆：",
        "  通常系統會自動提示允許",
    ],
    [  # Linux
        "systemd 方式（推薦）：",
        "$ sudo systemctl edit ollama",
        "",
        "加入：",
        "  [Service]",
        "  Environment=",
        "    \"OLLAMA_HOST=0.0.0.0\"",
        "",
        "$ sudo systemctl daemon-reload",
        "$ sudo systemctl restart ollama",
        "",
        "防火牆：",
        "$ sudo ufw allow 11434/tcp",
    ],
    note="設定後驗證：curl http://<SERVER_IP>:11434/api/tags"
)

# ── Slide 30: Section - Commands ──
make_section_slide(10, "常用指令速查表")

# ── Slide 26: Command reference ──
make_table_slide("10. 設定與診斷指令",
    ["指令", "說明"],
    [
        ["ai-review config show", "顯示所有設定"],
        ["ai-review config get <section> <key>", "取得單一設定值"],
        ["ai-review config set <section> <key> <value>", "設定值"],
        ["ai-review config init-template", "初始化 commit 模板"],
        ["ai-review config init-generate-prompt", "初始化 AI 生成 prompt"],
        ["ai-review health-check", "測試 LLM 連線"],
        ["ai-review -v", "Debug 模式審查"],
        ["ai-review --format markdown", "Markdown 輸出（貼 Issue）"],
        ["ai-review --format json", "JSON 輸出（CI/CD 整合）"],
    ]
)

# ── Slide 27: Hook commands ──
make_table_slide("10. Hook 管理指令",
    ["指令", "說明"],
    [
        ["ai-review hook install --template", "安裝 template hooks（推薦）"],
        ["ai-review hook install --global", "安裝 global hooks"],
        ["ai-review hook uninstall --template", "移除 template hooks"],
        ["ai-review hook status", "查看 hook 安裝狀態"],
        ["ai-review hook enable", "啟用當前 repo（自動安裝 hooks）"],
        ["hook enable --path <dir>", "啟用指定 repo（可重複）"],
        ["hook enable --all <dir>", "批量啟用目錄下所有 repo"],
        ["hook enable --all <dir> --list", "預覽所有 repo 狀態"],
        ["hook disable --all <dir>", "批量停用"],
    ]
)

# ── Slide 28: Section - Troubleshooting ──
make_section_slide(11, "疑難排解")

# ── Slide 29: Common issues ──
make_table_slide("11. 全平台通用問題",
    ["問題", "解法"],
    [
        ["command not found", "確認 pip Scripts 目錄在 PATH 中"],
        ["health-check 失敗", "確認 Ollama 服務：ollama list"],
        ["LLM 回應太慢", "換小模型：config set ollama model llama3.2"],
        ["Timeout 錯誤", "加大：config set ollama timeout 300"],
        ["Hook 不觸發", "確認：hook status + hook enable"],
        ["格式被擋", "格式須為 [tag] description"],
    ]
)

# ── Slide 31: Platform-specific issues ──
make_3col_slide("11. 平台專屬問題",
    [  # Windows
        "單引號 commit 失敗",
        "→ 用雙引號",
        "",
        "Hook 找不到 ai-review",
        "→ 用完整路徑：",
        "  /c/Users/<USER>/anaconda3",
        "  /Scripts/ai-review.exe",
        "",
        "中文亂碼",
        "→ chcp 65001",
        "",
        "pip install 檔案被鎖",
        "→ 關閉所有終端後重試",
    ],
    [  # macOS
        "python3 找不到",
        "→ brew install python@3.12",
        "",
        "Ollama 未啟動",
        "→ ollama serve &",
        "  或從應用程式啟動",
    ],
    [  # Linux
        "Ollama 服務未運行",
        "→ sudo systemctl start ollama",
        "",
        "權限不足",
        "→ pip install --user .",
        "  或使用 venv",
        "",
        "Hook 無執行權限",
        "→ chmod +x .git/hooks/*",
    ],
)

# ── Slide 32: Section - Uninstall ──
make_section_slide(12, "完整移除")

# ── Slide 33: Uninstall ──
make_3col_slide("12. 完整移除",
    [  # Windows
        "移除 hooks：",
        "$ ai-review hook uninstall \\",
        "    --template",
        "",
        "移除套件：",
        "$ pip uninstall ai-code-review",
        "",
        "移除設定：",
        "$ rmdir /s /q",
        "  %USERPROFILE%\\.config\\",
        "  ai-code-review",
        "",
        "移除 Ollama：",
        "  設定 > 應用程式 > 解除安裝",
    ],
    [  # macOS
        "移除 hooks：",
        "$ ai-review hook uninstall \\",
        "    --template",
        "",
        "移除套件：",
        "$ pip3 uninstall ai-code-review",
        "",
        "移除設定：",
        "$ rm -rf ~/.config/ai-code-review",
        "",
        "移除 Ollama：",
        "$ brew uninstall ollama",
    ],
    [  # Linux
        "移除 hooks：",
        "$ ai-review hook uninstall \\",
        "    --template",
        "",
        "移除套件：",
        "$ pip3 uninstall ai-code-review",
        "",
        "移除設定：",
        "$ rm -rf ~/.config/ai-code-review",
        "",
        "移除 Ollama：",
        "$ sudo systemctl stop ollama",
        "$ sudo systemctl disable ollama",
        "$ sudo rm /usr/local/bin/ollama",
    ],
)

# ── Appendix: Scenario - Multi repo (advanced) ──
make_content_slide("8. 情境二：管理多個 Repo — 進階設定（參考用）", [
    ("啟用指定 repo", {"size": 18, "bold": True, "color": ACCENT2}),
    ("$ ai-review hook enable --path /workspace/camera-hal \\", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("                        --path /workspace/kernel", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("  可用 --path 重複指定多個 repo", {"size": 13, "color": GRAY, "indent": 1}),
    ("", {"size": 8}),
    ("停用指定 repo", {"size": 18, "bold": True, "color": ACCENT3}),
    ("$ ai-review hook disable --path /workspace/kernel", {"size": 14, "color": ACCENT3, "indent": 1}),
    ("", {"size": 8}),
    ("批量停用", {"size": 18, "bold": True, "color": WIN_CLR}),
    ("$ ai-review hook disable --all /workspace", {"size": 14, "color": WIN_CLR, "indent": 1}),
    ("  停用目錄下所有 repo 的 ai-review hooks", {"size": 13, "color": GRAY, "indent": 1}),
])

# ── Appendix: Scenario - Shared LLM ──
make_content_slide("8. 情境三：團隊共用 LLM Server（參考用）", [
    ("架構", {"size": 18, "bold": True, "color": ACCENT}),
    ("  開發機 A (Windows) ──→", {"size": 15, "color": WIN_CLR}),
    ("  開發機 B (macOS)   ──→  LLM Server (192.168.1.100)", {"size": 15, "color": MAC_CLR}),
    ("  開發機 C (Linux)   ──→  Ollama + GPU", {"size": 15, "color": LNX_CLR}),
    ("", {"size": 10}),
    ("Server 端（一次）", {"size": 18, "bold": True, "color": ACCENT3}),
    ("$ sudo systemctl edit ollama", {"size": 14, "color": ACCENT3, "indent": 1}),
    ("  加入 Environment=\"OLLAMA_HOST=0.0.0.0\"", {"size": 13, "color": GRAY, "indent": 1}),
    ("$ sudo systemctl restart ollama", {"size": 14, "color": ACCENT3, "indent": 1}),
    ("", {"size": 6}),
    ("每台開發機（一次）", {"size": 18, "bold": True, "color": ACCENT2}),
    ("$ ai-review config set ollama base_url http://192.168.1.100:11434", {"size": 14, "color": ACCENT2, "indent": 1}),
    ("$ ai-review health-check", {"size": 14, "color": ACCENT2, "indent": 1}),
])

# ── End ──
make_title_slide(
    "Thank You",
    "Questions?  →  github.com/jame472518-design/ai-code-review"
)

# ── Save ──
import os
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "SOP-cross-platform.pptx")
prs.save(out_path)
print(f"PPT saved to: {out_path}")
